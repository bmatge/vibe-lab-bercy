#!/usr/bin/env python3
"""One-shot script: convert all DOCX/PPTX in docs/ to Markdown files."""

import os
import re
from docx import Document
from pptx import Presentation

DOCS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs')


def docx_to_md(path):
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append('')
            continue
        style = (para.style.name if para.style and para.style.name else '').lower()
        if 'heading 1' in style or style == 'title':
            lines.append(f'# {text}')
        elif 'heading 2' in style:
            lines.append(f'## {text}')
        elif 'heading 3' in style:
            lines.append(f'### {text}')
        elif 'heading 4' in style:
            lines.append(f'#### {text}')
        elif 'list' in style:
            lines.append(f'- {text}')
        else:
            lines.append(text)

    # Process tables
    for table in doc.tables:
        lines.append('')
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        lines.append('| ' + ' | '.join(headers) + ' |')
        lines.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')
        for row in table.rows[1:]:
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            lines.append('| ' + ' | '.join(cells) + ' |')
        lines.append('')

    md = '\n'.join(lines)
    # Collapse 3+ blank lines into 2
    md = re.sub(r'\n{3,}', '\n\n', md)
    return md.strip() + '\n'


def pptx_to_md(path):
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        title = ''
        body_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                is_title = (shape == slide.shapes.title)
                if not is_title:
                    try:
                        is_title = shape.placeholder_format.idx == 0
                    except (ValueError, AttributeError):
                        pass
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if is_title:
                        title = text
                    else:
                        body_parts.append(text)
        if title:
            lines.append(f'## Slide {i} — {title}')
        else:
            lines.append(f'## Slide {i}')
        lines.append('')
        for part in body_parts:
            lines.append(part)
            lines.append('')

    md = '\n'.join(lines)
    md = re.sub(r'\n{3,}', '\n\n', md)
    return md.strip() + '\n'


def main():
    for filename in sorted(os.listdir(DOCS_DIR)):
        filepath = os.path.join(DOCS_DIR, filename)
        basename, ext = os.path.splitext(filename)
        ext = ext.lower()

        if ext == '.docx':
            print(f'Converting {filename} ...')
            md = docx_to_md(filepath)
            out = os.path.join(DOCS_DIR, f'{basename}.md')
            with open(out, 'w', encoding='utf-8') as f:
                f.write(md)
            print(f'  → {basename}.md ({len(md)} chars)')

        elif ext == '.pptx':
            print(f'Converting {filename} ...')
            md = pptx_to_md(filepath)
            out = os.path.join(DOCS_DIR, f'{basename}.md')
            with open(out, 'w', encoding='utf-8') as f:
                f.write(md)
            print(f'  → {basename}.md ({len(md)} chars)')


if __name__ == '__main__':
    main()
